"""
build_slide.py — '시냅스 최종 공식 + 8파라미터를 찾는 6단계' 슬라이드(.pptx) 생성
상단: 색깔별 공식 / 하단: 6단계 그림 카드(번호 배지 + 파라미터 색 매칭).
실행: conda activate ca1sim
      python papers/01_Ecker2020_CA1_synaptic/slides/build_slide.py
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

THIS = os.path.dirname(os.path.abspath(__file__))
PAPER = os.path.dirname(THIS)
FIG = os.path.join(PAPER, "03_synapses", "figures")
OUT = os.path.join(THIS, "synapse_model_8params_6steps.pptx")

FONT = "Malgun Gothic"
NAVY = RGBColor(0x1E, 0x27, 0x61)
BLACK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x78, 0x78, 0x78)
BLUE = RGBColor(0x15, 0x65, 0xC0)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE0, 0x7B, 0x00)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_font(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs = [(text, size, bold, color)] → 한 단락에 색깔별 run."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    for text, size, bold, color in runs:
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        set_font(r)
    return tb


def fit(img_path, box_w, box_h):
    """이미지를 (box_w, box_h) 안에 비율 유지로 맞춘 (w,h) inch 반환."""
    with Image.open(img_path) as im:
        a = im.width / im.height
    if box_w / box_h > a:        # 높이 제한
        h = box_h; w = box_h * a
    else:                        # 너비 제한
        w = box_w; h = box_w / a
    return w, h


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 제목
    add_text(slide, 0.45, 0.22, 12.4, 0.7,
             [("CA1 시냅스 모델 — 최종 공식 & 8개 파라미터를 찾는 6단계", 30, True, NAVY)])

    # 공식 (색깔별 파라미터)
    f = [("I(t) = ", 23, False, BLACK), ("ĝ", 25, True, PURPLE),
         (" · biexp(t; ", 23, False, BLACK), ("τ_rise, τ_decay", 25, True, BLUE),
         (") · TM(", 23, False, BLACK), ("U_SE, D, F", 25, True, GREEN),
         (", ", 23, False, BLACK), ("N_RRP", 25, True, ORANGE),
         (") · (Vm − ", 23, False, BLACK), ("E_rev", 25, True, BLUE), (")", 23, False, BLACK)]
    add_text(slide, 0.45, 1.05, 12.4, 0.7, f, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 공식 범례
    leg = [("파라미터 색 = 찾는 단계:  ", 12, False, GRAY), ("biexp(단계3)", 12, True, BLUE),
           (" · ", 12, False, GRAY), ("TM(단계4)", 12, True, GREEN), (" · ", 12, False, GRAY),
           ("확률방출(단계5)", 12, True, ORANGE), (" · ", 12, False, GRAY), ("ĝ 보정(단계6)", 12, True, PURPLE)]
    add_text(slide, 0.45, 1.72, 12.4, 0.35, leg, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 6단계 카드
    steps = [
        (1, "단계1 · 해부", "축삭-수상돌기 분포", "—", GRAY, "1_innervation.png"),
        (2, "단계2 · 해부", "연결당 시냅스 수", "—", GRAY, "2_num_synapses.png"),
        (3, "단계3 · §2.3", "biexp 전도도·전류", "E_rev, τ_rise, τ_decay", BLUE, "s1_biexp_conductance.png"),
        (4, "단계4 · §2.4", "TM 단기가소성", "U_SE, D, F", GREEN, "s2_tm_stp.png"),
        (5, "단계5 · §2.5", "확률 다소포 방출", "N_RRP", ORANGE, "s3_stochastic_mvr.png"),
        (6, "단계6 · §2.6", "ĝ 보정", "ĝ", PURPLE, "s1_calibrate_ghat.png"),
    ]
    mx, gap = 0.4, 0.25
    col_w = (13.333 - 2 * mx - 2 * gap) / 3
    cols = [mx + i * (col_w + gap) for i in range(3)]
    rows = [2.25, 4.88]
    card_h = 2.45

    for idx, (n, lab, desc, params, color, fname) in enumerate(steps):
        cx = cols[idx % 3]; cy = rows[idx // 3]
        # 번호 배지 (색 원 + 숫자)
        d = 0.42
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(d), Inches(d))
        badge.fill.solid(); badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        btf = badge.text_frame; btf.word_wrap = False
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = str(n); br.font.size = Pt(16); br.font.bold = True
        br.font.color.rgb = WHITE; set_font(br)
        # 라벨/설명/파라미터
        add_text(slide, cx + 0.5, cy - 0.04, col_w - 0.5, 0.55,
                 [(f"{lab}  {desc}\n", 12.5, True, color),
                  (f"찾는 파라미터: {params}", 11, True, color) if color != GRAY
                  else ("해부 검증 (파라미터 아님)", 11, False, GRAY)],
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        # 그림 (비율 유지 + 중앙)
        bw, bh = col_w, card_h - 0.62
        w, h = fit(os.path.join(FIG, fname), bw, bh)
        ix = cx + (bw - w) / 2; iy = cy + 0.58 + (bh - h) / 2
        slide.shapes.add_picture(os.path.join(FIG, fname), Inches(ix), Inches(iy),
                                 Inches(w), Inches(h))

    prs.save(OUT)
    print(f"[pptx] {OUT}")


if __name__ == "__main__":
    main()
