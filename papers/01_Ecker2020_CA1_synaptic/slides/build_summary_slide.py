"""
build_summary_slide.py — 1장 요약 슬라이드(.pptx)
  ① 시냅스 모델 + 8파라미터·6단계 (slides/preview.png)
  ② 9클래스 단기가소성(STP) 재현·검증 (03_synapses/figures/4-1_reproduce_fig5.png)
  ③ 축소 CA1 마이크로서킷 (04_network/figures/1_connectivity.png)
실행: <ca1sim python> papers/01_Ecker2020_CA1_synaptic/slides/build_summary_slide.py
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

THIS = os.path.dirname(os.path.abspath(__file__))
PAPER = os.path.dirname(THIS)
OUT = os.path.join(THIS, "CA1_summary_1slide.pptx")

FONT = "Malgun Gothic"
NAVY = RGBColor(0x1E, 0x27, 0x61)
BLACK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x70, 0x70, 0x70)
BLUE = RGBColor(0x15, 0x65, 0xC0)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

IMG = {
    "model": os.path.join(THIS, "preview.png"),
    "stp":   os.path.join(PAPER, "03_synapses", "figures", "4-1_reproduce_fig5.png"),
    "net":   os.path.join(PAPER, "04_network", "figures", "1_connectivity.png"),
}


def set_font(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    for text, size, bold, color in runs:
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        set_font(r)
    return tb


def fit(img_path, box_w, box_h):
    with Image.open(img_path) as im:
        a = im.width / im.height
    if box_w / box_h > a:
        h = box_h; w = box_h * a
    else:
        w = box_w; h = box_w / a
    return w, h


def badge(slide, x, y, n, color, d=0.38):
    b = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
    tf = b.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n); r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = WHITE; set_font(r)


def panel(slide, n, color, caption, img_key, box):
    """번호 배지 + 캡션 + 이미지(비율 유지·중앙)."""
    bx, by, bw, bh = box
    badge(slide, bx, by, n, color)
    add_text(slide, bx + 0.46, by - 0.02, bw - 0.46, 0.42,
             [(caption, 13, True, color)], anchor=MSO_ANCHOR.MIDDLE)
    iy0 = by + 0.5
    ibh = bh - 0.5
    w, h = fit(IMG[img_key], bw, ibh)
    ix = bx + (bw - w) / 2; iy = iy0 + (ibh - h) / 2
    slide.shapes.add_picture(IMG[img_key], Inches(ix), Inches(iy), Inches(w), Inches(h))


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    add_text(slide, 0.45, 0.2, 12.45, 0.62,
             [("CA1 in silico — 시냅스 모델 → 9클래스 검증 → 마이크로서킷",
               24, True, NAVY)], anchor=MSO_ANCHOR.MIDDLE)

    # 상단 좌: 시냅스 모델 + 6단계 / 상단 우: 9클래스 STP 재현
    panel(slide, 1, NAVY,  "시냅스 모델 + 8 파라미터 · 6단계",
          "model", (0.45, 0.95, 6.15, 3.55))
    panel(slide, 2, BLUE,  "9클래스 단기가소성(STP) 재현 · 검증",
          "stp",   (6.95, 0.95, 5.95, 3.55))
    # 하단: 마이크로서킷(와이드)
    panel(slide, 3, GREEN, "축소 CA1 마이크로서킷 (122세포 · 553연결 · 9클래스 시냅스)",
          "net",   (0.45, 4.75, 12.45, 2.55))

    prs.save(OUT)
    print(f"[pptx] {OUT}")


if __name__ == "__main__":
    main()
