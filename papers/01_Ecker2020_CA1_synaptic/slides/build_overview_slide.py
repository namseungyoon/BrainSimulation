"""
build_overview_slide.py — 1장 '방법 개요' 슬라이드(.pptx)
상단: ① 뉴런 4종(형태) · ② 9 경로 클래스 · (오른쪽) ④ 마이크로서킷 계획(상세)
하단: ③ 시냅스 파라미터 6단계 절차 (1→6 가로 나열)
실행: <ca1sim python> papers/01_Ecker2020_CA1_synaptic/slides/build_overview_slide.py
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

THIS = os.path.dirname(os.path.abspath(__file__))
PAPER = os.path.dirname(THIS)
SFIG = os.path.join(PAPER, "03_synapses", "figures")
NFIG = os.path.join(PAPER, "04_network", "figures")
OUT = os.path.join(THIS, "CA1_overview_1slide.pptx")

FONT = "Malgun Gothic"
NAVY = RGBColor(0x1E, 0x27, 0x61)
BLACK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x70, 0x70, 0x70)
BLUE = RGBColor(0x15, 0x65, 0xC0)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE0, 0x7B, 0x00)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TINT = RGBColor(0xF2, 0xF4, 0xF9)


def set_font(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04); tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = align
    for text, size, bold, color in runs:
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; set_font(r)
    return tb


def add_bullets(slide, x, y, w, h, lines, size=10.0, color=BLACK):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06); tf.margin_top = Inches(0.04)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run(); r.text = "• " + ln
        r.font.size = Pt(size); r.font.color.rgb = color; set_font(r)
    return tb


def fit(img_path, box_w, box_h):
    with Image.open(img_path) as im:
        a = im.width / im.height
    if box_w / box_h > a:
        h = box_h; w = box_h * a
    else:
        w = box_w; h = box_w / a
    return w, h


def picture(slide, path, bx, by, bw, bh, valign="center"):
    w, h = fit(path, bw, bh)
    iy = by if valign == "top" else by + (bh - h) / 2
    slide.shapes.add_picture(path, Inches(bx + (bw - w) / 2), Inches(iy), Inches(w), Inches(h))


def badge(slide, x, y, n, color, d=0.36):
    b = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
    tf = b.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n); r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = WHITE; set_font(r)


def panel_card(slide, x, y, w, h):
    c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = TINT; c.line.fill.background()
    c.shadow.inherit = False
    return c


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, 0.4, 0.12, 12.5, 0.5,
             [("CA1 in silico — 마이크로서킷 구축 방법 개요", 24, True, NAVY)],
             anchor=MSO_ANCHOR.MIDDLE)

    # ── 상단 3분할 ─────────────────────────────────────────────
    top_y, top_h = 0.72, 2.95
    # ① 뉴런 4종 형태
    add_text(slide, 0.35, top_y, 4.6, 0.34, [("① 뉴런 4종 (e-type 대표) 형태", 13, True, NAVY)])
    picture(slide, os.path.join(NFIG, "celltypes_4.png"), 0.35, top_y + 0.36, 4.6, top_h - 0.36)
    # ② 9 경로 클래스
    add_text(slide, 5.1, top_y, 3.7, 0.34, [("② 9 경로 클래스 (Table 3) · STP 재현", 13, True, BLUE)])
    picture(slide, os.path.join(SFIG, "4-1_reproduce_fig5.png"), 5.1, top_y + 0.36, 3.7, top_h - 0.36)
    # ④ 마이크로서킷 계획 (오른쪽, 카드 + 글머리)
    px, pw = 8.95, 4.05
    panel_card(slide, px, top_y - 0.02, pw, top_h + 0.05)
    add_text(slide, px + 0.05, top_y + 0.02, pw - 0.1, 0.34,
             [("④ 축소 마이크로서킷 계획", 13, True, GREEN)])
    add_bullets(slide, px + 0.05, top_y + 0.4, pw - 0.1, top_h - 0.4, [
        "규모: 122세포 (PC 100 · PV 10 · cAC 6 · bAC 6)",
        "배치: 3D 층 SP/SO/SR (500µm 타일)",
        "연결: 거리의존 확률(정점×Gaussian, 공간상수 150µm) → 553연결",
        "시냅스: 9클래스(Table 3) 자동 배정",
        "방출: EMS 확률 시냅스 → 고정 dt 실행",
        "구동: PC 외부 입력 → raster · e-type별 발화율",
        "다음(세션 D): paired recording — 연결별 PSP 분포",
    ])

    # ── 하단: 6단계 절차 (가로 1→6) ────────────────────────────
    add_text(slide, 0.35, 3.78, 12.6, 0.34,
             [("③ 시냅스 파라미터 8개를 정하는 6단계 절차  ", 13, True, BLACK),
              ("(단계3 E_rev·τ · 단계4 U,D,F · 단계5 N_RRP · 단계6 ĝ)", 11, False, GRAY)])
    steps = [
        (1, "축삭-수상돌기 분포", GRAY, os.path.join(SFIG, "1_innervation.png")),
        (2, "연결당 시냅스 수", GRAY, os.path.join(SFIG, "2_num_synapses.png")),
        (3, "biexp 전도도", BLUE, os.path.join(SFIG, "3_biexp_conductance.png")),
        (4, "TM 단기가소성", GREEN, os.path.join(SFIG, "4_class_PC-PC_E2.png")),
        (5, "확률 다소포 방출", ORANGE, os.path.join(SFIG, "5_stochastic_mvr.png")),
        (6, "ĝ 보정", PURPLE, os.path.join(SFIG, "6_calibrate_ghat.png")),
    ]
    mx, gap, n = 0.35, 0.14, 6
    col_w = (13.333 - 2 * mx - (n - 1) * gap) / n
    fy = 4.60
    for i, (num, lab, color, path) in enumerate(steps):
        cx = mx + i * (col_w + gap)
        badge(slide, cx, 4.18, num, color)
        add_text(slide, cx + 0.4, 4.16, col_w - 0.4, 0.42, [(lab, 10.5, True, color)],
                 anchor=MSO_ANCHOR.MIDDLE)
        picture(slide, path, cx, fy, col_w, 7.34 - fy, valign="top")

    prs.save(OUT)
    print(f"[pptx] {OUT}")


if __name__ == "__main__":
    main()
